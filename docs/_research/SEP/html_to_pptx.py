#!/usr/bin/env python3
from __future__ import annotations

import argparse
import datetime as _dt
import shutil
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


def _file_url(path: Path) -> str:
    return path.resolve().as_uri()


def _timestamp() -> str:
    return _dt.datetime.now().strftime("%Y%m%d-%H%M%S")


def _backup_existing(path: Path) -> None:
    if not path.exists():
        return
    backup = path.with_suffix(f".bak-{_timestamp()}{path.suffix}")
    shutil.move(str(path), str(backup))


def _wait_images_loaded(page) -> None:
    page.evaluate(
        """
        () => Promise.all(
          Array.from(document.images).map(img => {
            if (img.complete) return true;
            return new Promise(resolve => {
              const done = () => resolve(true);
              img.addEventListener('load', done, { once: true });
              img.addEventListener('error', done, { once: true });
            });
          })
        )
        """
    )


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Convert an HTML slide deck (with .slide elements) into a PPTX by screenshotting each slide."
    )
    parser.add_argument("html", type=Path, help="Path to the HTML deck")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Output .pptx path (default: same name as input, .pptx)",
    )
    parser.add_argument("--width", type=int, default=1920, help="Viewport width (px)")
    parser.add_argument("--height", type=int, default=1080, help="Viewport height (px)")
    args = parser.parse_args()

    html_path: Path = args.html
    if not html_path.exists():
        raise SystemExit(f"HTML not found: {html_path}")

    out_path: Path = args.output or html_path.with_suffix(".pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="html_to_pptx_") as tmpdir:
        tmpdir_path = Path(tmpdir)

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": args.width, "height": args.height})
            page.goto(_file_url(html_path))
            page.wait_for_load_state("domcontentloaded")

            page.add_style_tag(
                content="""
                .nav, .counter { display: none !important; }
                * , *::before, *::after { transition: none !important; animation: none !important; }
                """
            )
            _wait_images_loaded(page)

            slide_count: int = page.evaluate("() => document.querySelectorAll('.slide').length")
            if slide_count <= 0:
                raise SystemExit("No slides found (expected elements with class .slide).")

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]

            for i in range(slide_count):
                page.evaluate(
                    """
                    (i) => {
                      const slides = Array.from(document.querySelectorAll('.slide'));
                      slides.forEach((s, j) => {
                        s.classList.remove('active', 'prev');
                        if (j === i) s.classList.add('active');
                        else if (j < i) s.classList.add('prev');
                      });
                      const cur = document.getElementById('cur');
                      if (cur) cur.textContent = String(i + 1);
                    }
                    """,
                    i,
                )
                page.wait_for_timeout(30)
                _wait_images_loaded(page)

                img_path = tmpdir_path / f"slide-{i+1:02d}.png"
                page.screenshot(path=str(img_path), full_page=False)

                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(
                    str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height
                )

            browser.close()

        _backup_existing(out_path)
        prs.save(str(out_path))

    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

