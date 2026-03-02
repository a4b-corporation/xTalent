#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "bg": RGBColor(0x0A, 0x0E, 0x1A),
    "surface": RGBColor(0x11, 0x18, 0x27),
    "card": RGBColor(0x1A, 0x22, 0x35),
    "text": RGBColor(0xE2, 0xE8, 0xF0),
    "dim": RGBColor(0x94, 0xA3, 0xB8),
    "purple": RGBColor(0x63, 0x66, 0xF1),
    "cyan": RGBColor(0x06, 0xB6, 0xD4),
    "amber": RGBColor(0xF5, 0x9E, 0x0B),
    "green": RGBColor(0x10, 0xB9, 0x81),
    "pink": RGBColor(0xEC, 0x48, 0x99),
}


def _safe_text(text: str) -> str:
    t = re.sub(r"\s+", " ", text or "").strip()
    return t


def _img_contain(path: Path, max_w_in: float, max_h_in: float) -> tuple[float, float]:
    with Image.open(path) as im:
        w, h = im.size
    if w <= 0 or h <= 0:
        return max_w_in, max_h_in
    scale = min(max_w_in / (w / 96.0), max_h_in / (h / 96.0))
    return (w / 96.0) * scale, (h / 96.0) * scale


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_tag(slide, text: str) -> None:
    if not text:
        return
    x, y, w, h = Inches(0.6), Inches(0.4), Inches(3.0), Inches(0.45)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["surface"]
    shape.line.color.rgb = COLORS["purple"]
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = COLORS["purple"]
    p.alignment = PP_ALIGN.CENTER


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    title = _safe_text(title)
    subtitle = _safe_text(subtitle or "")
    if title:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(1.1))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.0))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["dim"]
        p.alignment = PP_ALIGN.CENTER


@dataclass(frozen=True)
class Card:
    title: str
    body: str
    accent: str | None = None


def _add_cards_grid(
    slide,
    cards: list[Card],
    cols: int,
    *,
    x0: float = 0.7,
    y0: float = 2.6,
    w: float = 11.9,
    h: float = 4.6,
) -> None:
    if not cards:
        return
    gap = 0.25
    rows = (len(cards) + cols - 1) // cols
    card_w = (w - gap * (cols - 1)) / cols
    card_h = (h - gap * (rows - 1)) / max(rows, 1)

    for idx, card in enumerate(cards):
        r = idx // cols
        c = idx % cols
        x = x0 + c * (card_w + gap)
        y = y0 + r * (card_h + gap)

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(card_w),
            Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["card"]
        shape.line.color.rgb = RGBColor(0x22, 0x2B, 0x3D)
        shape.line.width = Pt(1)

        if card.accent and card.accent in COLORS:
            stripe = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(0.08),
                Inches(card_h),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = COLORS[card.accent]
            stripe.line.fill.background()

        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = _safe_text(card.title)
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLORS["text"]

        if card.body:
            p2 = tf.add_paragraph()
            p2.text = _safe_text(card.body)
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLORS["dim"]


def _add_bullets(
    slide,
    title: str,
    items: list[str],
    *,
    x: float = 1.2,
    y: float = 2.2,
    w: float = 11.0,
    h: float = 4.9,
) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()

    if title:
        p = tf.paragraphs[0]
        p.text = _safe_text(title)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]
    for i, it in enumerate(items):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = _safe_text(it)
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["dim"]


def _add_image_center(slide, img_path: Path, *, y_top: float = 2.1, max_h: float = 4.9) -> None:
    if not img_path.exists():
        return
    max_w = 12.2
    w_in, h_in = _img_contain(img_path, max_w, max_h)
    x = (13.333 - w_in) / 2.0
    y = y_top + max(0.0, (max_h - h_in) / 2.0)
    slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(w_in), height=Inches(h_in))


def _parse_cards(slide_div) -> list[Card]:
    cards: list[Card] = []
    for card_div in slide_div.select(".card"):
        title = _safe_text(card_div.select_one("h3").get_text(" ", strip=True) if card_div.select_one("h3") else "")
        body = _safe_text(card_div.get_text(" ", strip=True))
        if title and body.startswith(title):
            body = _safe_text(body[len(title) :])
        accent = None
        classes = set(card_div.get("class", []))
        if "purple-border" in classes:
            accent = "purple"
        elif "cyan-border" in classes:
            accent = "cyan"
        elif "amber-border" in classes:
            accent = "amber"
        elif "green-border" in classes:
            accent = "green"
        cards.append(Card(title=title, body=body, accent=accent))
    return cards


def _find_slide_image(slide_div) -> str | None:
    img = slide_div.select_one(".img-container img") or slide_div.select_one("img")
    if not img:
        return None
    src = (img.get("src") or "").strip()
    return src or None


def _build_pptx(html_path: Path, out_path: Path) -> None:
    html = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "lxml")
    slide_divs = soup.select("div.slide")
    if not slide_divs:
        raise SystemExit("No slides found (expected div.slide).")

    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    for slide_div in slide_divs:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, COLORS["bg"])

        tag = _safe_text(slide_div.select_one(".tag").get_text(" ", strip=True) if slide_div.select_one(".tag") else "")
        h1 = _safe_text(slide_div.select_one("h1").get_text(" ", strip=True) if slide_div.select_one("h1") else "")
        h2 = _safe_text(slide_div.select_one("h2").get_text(" ", strip=True) if slide_div.select_one("h2") else "")
        subtitle = _safe_text(slide_div.select_one(".subtitle").get_text(" ", strip=True) if slide_div.select_one(".subtitle") else "")

        if tag:
            _add_tag(slide, tag)

        title = h1 or h2
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(0.9))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(34 if h2 else 54)
            p.font.bold = True
            p.font.color.rgb = COLORS["text"]
            p.alignment = PP_ALIGN.CENTER

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(0.9))
            tf = sub_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS["dim"]
            p.alignment = PP_ALIGN.CENTER

        img_src = _find_slide_image(slide_div)
        if img_src and img_src.startswith("assets/"):
            _add_image_center(slide, (html_path.parent / img_src).resolve(), y_top=2.4, max_h=4.7)
            continue

        grid3 = slide_div.select_one(".grid3")
        grid4 = slide_div.select_one(".grid4")
        grid5 = slide_div.select_one(".grid5")
        grid2 = slide_div.select_one(".grid2")

        if grid4 and slide_div.select_one(".stat"):
            stats = []
            for stat in slide_div.select(".stat"):
                num = _safe_text(stat.select_one(".num").get_text(" ", strip=True) if stat.select_one(".num") else "")
                label = _safe_text(
                    stat.select_one(".label").get_text(" ", strip=True) if stat.select_one(".label") else ""
                )
                if num or label:
                    stats.append(Card(title=num, body=label, accent="cyan"))
            _add_cards_grid(slide, stats, cols=4, x0=0.7, y0=2.7, w=11.9, h=2.6)
            continue

        if slide_div.select_one(".flow-step"):
            items = []
            for step in slide_div.select(".flow-step"):
                txt = _safe_text(step.get_text(" ", strip=True))
                if txt:
                    items.append(txt)
            _add_bullets(slide, "", items, x=1.4, y=2.3, w=10.6, h=4.9)
            continue

        if slide_div.select_one(".phase-card"):
            cards: list[Card] = []
            for phase in slide_div.select(".phase-card"):
                phase_title = _safe_text(phase.select_one("h3").get_text(" ", strip=True) if phase.select_one("h3") else "")
                period = _safe_text(phase.select_one(".period").get_text(" ", strip=True) if phase.select_one(".period") else "")
                lis = [_safe_text(li.get_text(" ", strip=True)) for li in phase.select("li")]
                body = "\n".join([period] + lis if period else lis)
                accent = "purple"
                classes = set(phase.get("class", []))
                if "p1" in classes:
                    accent = "purple"
                elif "p2" in classes:
                    accent = "cyan"
                elif "p3" in classes:
                    accent = "green"
                cards.append(Card(title=phase_title, body=body, accent=accent))
            _add_cards_grid(slide, cards, cols=3, x0=0.7, y0=2.6, w=11.9, h=4.6)
            continue

        cards = _parse_cards(slide_div)
        if cards and (grid5 or grid4 or grid3 or grid2):
            cols = 3 if grid3 else 5 if grid5 else 4 if grid4 else 2
            _add_cards_grid(slide, cards, cols=cols, x0=0.7, y0=2.6, w=11.9, h=4.6)
            continue

        if slide_div.select_one(".cta-box"):
            _set_bg(slide, COLORS["surface"])
            h1_cta = _safe_text(slide_div.select_one(".cta-box h1").get_text(" ", strip=True))
            p_cta = _safe_text(slide_div.select_one(".cta-box p").get_text(" ", strip=True))
            if h1_cta:
                tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.6), Inches(2.0))
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = h1_cta
                p.font.size = Pt(34)
                p.font.bold = True
                p.font.color.rgb = COLORS["text"]
                p.alignment = PP_ALIGN.CENTER
            if p_cta:
                tb = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.9), Inches(1.2))
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = p_cta
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS["dim"]
                p.alignment = PP_ALIGN.CENTER

            pills = [_safe_text(s.get_text(" ", strip=True)) for s in slide_div.select(".cta-box span")]
            if pills:
                x = 1.2
                y = 5.2
                for pill in pills:
                    w = max(1.6, min(4.2, 0.12 * len(pill) + 1.0))
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.5)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = COLORS["card"]
                    shape.line.fill.background()
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = pill
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = COLORS["text"]
                    p.alignment = PP_ALIGN.CENTER
                    x += w + 0.2
                    if x > 11.5:
                        x = 1.2
                        y += 0.65
            continue

        fallback = _safe_text(slide_div.get_text("\n", strip=True))
        if fallback:
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.8))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = fallback
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS["dim"]

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert xTalent HTML slide deck to an editable PPTX (no screenshots).")
    parser.add_argument("html", type=Path)
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Output .pptx path (default: <html stem>-editable.pptx next to HTML)",
    )
    args = parser.parse_args()
    html_path: Path = args.html
    out_path = args.output or html_path.with_name(f"{html_path.stem}-editable.pptx")
    _build_pptx(html_path, out_path)
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

